import io
from typing import List, Dict, Any, Optional
from app.modules.data_output.exporters.base_exporter import BaseExporter

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
except ImportError:
    Presentation = None


class PPTXExporter(BaseExporter):
    """PPTX Presentation Exporter utilizing python-pptx."""

    async def export_query_results(
        self,
        columns: List[str],
        rows: List[List[Any]],
        title: str = "Query Results",
        options: Optional[Dict[str, Any]] = None,
    ) -> bytes:
        if Presentation is not None:
            prs = Presentation()
            slide_layout = prs.slide_layouts[5]  # Title & Blank content
            slide = prs.slides.add_slide(slide_layout)
            
            shapes = slide.shapes
            title_shape = shapes.title
            title_shape.text = title
            
            # Table layout
            max_rows = min(len(rows), 15)
            rows_cnt = max_rows + 1
            cols_cnt = len(columns)
            
            left, top, width, height = Inches(0.5), Inches(1.5), Inches(9.0), Inches(5.0)
            table_shape = shapes.add_table(rows_cnt, cols_cnt, left, top, width, height)
            table = table_shape.table
            
            for col_idx, col in enumerate(columns):
                table.cell(0, col_idx).text = col
                
            for row_idx in range(max_rows):
                for col_idx in range(cols_cnt):
                    val = rows[row_idx][col_idx] if col_idx < len(rows[row_idx]) else ""
                    table.cell(row_idx + 1, col_idx).text = str(val) if val is not None else ""
                    
            output = io.BytesIO()
            prs.save(output)
            return output.getvalue()
            
        # Fallback binary string if python-pptx is not installed
        return f"PPTX export for '{title}' with {len(rows)} rows.".encode("utf-8")

    async def export_report(
        self,
        title: str,
        content_structure: Dict[str, Any],
        options: Optional[Dict[str, Any]] = None,
    ) -> bytes:
        if Presentation is not None:
            prs = Presentation()
            # Title Slide
            title_slide_layout = prs.slide_layouts[0]
            slide = prs.slides.add_slide(title_slide_layout)
            slide.shapes.title.text = title
            
            sections = content_structure.get("sections", [])
            for sec in sections:
                bullet_layout = prs.slide_layouts[1]
                b_slide = prs.slides.add_slide(bullet_layout)
                b_slide.shapes.title.text = sec.get("title", "Section")
                tf = b_slide.placeholders[1].text_frame
                tf.text = sec.get("content", "")
                
            output = io.BytesIO()
            prs.save(output)
            return output.getvalue()
            
        return await self.export_query_results(columns=["Report Title"], rows=[[title]], title=title)
